#!/usr/bin/env python3
"""Yome PPT chart engine.

Runs as a subprocess from PowerPointBridge.swift (and the CLI dispatcher).
All I/O is JSON over argv[1] — the Swift side serialises a single request
object and we echo a single response object on stdout.

Request shape:
    {
      "op":    "charts" | "add" | "update" | "delete",
      "file":  "/abs/path/to/file.pptx",
      "slide": 1-based int,
      "shape": optional 1-based int (update/delete),
      "type":  "column" | "bar" | "line" | "pie" | "area" | "scatter",
      "title": optional str,
      "csv":   CSV string (add/update),
      "left":  pt float,
      "top":   pt float,
      "width": pt float,
      "height":pt float
    }

Response shape (always a single JSON line on stdout):
    {"ok": true,  "result": {...}}
    {"ok": false, "error": "message"}

Requires python-pptx. We surface a clear message if it is missing.
"""
from __future__ import annotations

import json
import sys
import traceback
from typing import Any, Dict, List


def _fail(msg: str) -> None:
    sys.stdout.write(json.dumps({"ok": False, "error": msg}))
    sys.stdout.flush()
    sys.exit(0)


try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
except Exception as e:  # pragma: no cover - env dependent
    _fail(
        "python-pptx is required for chart operations. "
        "Install it with:  pip3 install python-pptx   (error: " + str(e) + ")"
    )


_CHART_ALIAS = {
    "column":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar":     XL_CHART_TYPE.BAR_CLUSTERED,
    "line":    XL_CHART_TYPE.LINE,
    "pie":     XL_CHART_TYPE.PIE,
    "area":    XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


def _reverse_chart_type(xl_type: Any) -> str:
    for alias, val in _CHART_ALIAS.items():
        if val == xl_type:
            return alias
    return str(xl_type)


def _parse_csv(raw: str):
    text = raw.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        raise ValueError("empty CSV data")
    lines = [l for l in text.split("\n") if l.strip()]
    if len(lines) < 2:
        raise ValueError("CSV needs a header row + at least 1 data row")
    header = [c.strip() for c in lines[0].split(",")]
    rows = [[c.strip() for c in l.split(",")] for l in lines[1:]]
    series_names = header[1:]
    categories = [r[0] if r else "" for r in rows]
    series: List[List[float]] = []
    for si, _name in enumerate(series_names):
        col = []
        for r in rows:
            try:
                col.append(float(r[si + 1]))
            except (IndexError, ValueError):
                col.append(0.0)
        series.append(col)
    return categories, series_names, series


def _cm_from_pt(v):
    # python-pptx accepts Pt/Emu directly; keep pt input.
    return Pt(float(v))


def _chart_list(prs, slide_idx: int):
    if slide_idx < 1 or slide_idx > len(prs.slides):
        raise ValueError(f"slide {slide_idx} out of range")
    slide = prs.slides[slide_idx - 1]
    out = []
    for i, shape in enumerate(slide.shapes, start=1):
        if shape.has_chart:
            chart = shape.chart
            title = ""
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    title = chart.chart_title.text_frame.text
            except Exception:
                pass
            series_count = 0
            try:
                series_count = len(chart.series)
            except Exception:
                pass
            out.append({
                "shape": i,
                "type": _reverse_chart_type(chart.chart_type),
                "title": title,
                "series": series_count,
            })
    return out


def _build_category_data(title: str, csv: str) -> CategoryChartData:
    categories, series_names, series = _parse_csv(csv)
    data = CategoryChartData()
    data.categories = categories
    if not series_names:
        # Single anonymous series = treat the only value column
        series_names = ["Series 1"]
    for name, values in zip(series_names, series):
        data.add_series(name, values)
    return data


def _build_xy_data(csv: str) -> XyChartData:
    # For scatter: expect header X, Y1[, Y2...]; first column = x, rest = y series.
    header_line, *rest = [l for l in csv.replace("\\n", "\n").split("\n") if l.strip()]
    header = [c.strip() for c in header_line.split(",")]
    rows = [[c.strip() for c in l.split(",")] for l in rest]
    series_names = header[1:] or ["Series 1"]
    data = XyChartData()
    for si, name in enumerate(series_names):
        s = data.add_series(name)
        for r in rows:
            try:
                x = float(r[0])
                y = float(r[si + 1])
            except (IndexError, ValueError):
                continue
            s.add_data_point(x, y)
    return data


def op_add(req: Dict[str, Any]) -> Dict[str, Any]:
    path = req["file"]
    prs = Presentation(path)
    slide_idx = int(req["slide"])
    if slide_idx < 1 or slide_idx > len(prs.slides):
        raise ValueError(f"slide {slide_idx} out of range")
    slide = prs.slides[slide_idx - 1]

    chart_alias = (req.get("type") or "column").lower()
    xl_type = _CHART_ALIAS.get(chart_alias, XL_CHART_TYPE.COLUMN_CLUSTERED)

    left = _cm_from_pt(req.get("left", 80))
    top = _cm_from_pt(req.get("top", 120))
    width = _cm_from_pt(req.get("width", 480))
    height = _cm_from_pt(req.get("height", 300))

    if chart_alias == "scatter":
        data = _build_xy_data(req.get("csv", ""))
    else:
        data = _build_category_data(req.get("title") or "", req.get("csv", ""))

    graphic_frame = slide.shapes.add_chart(xl_type, left, top, width, height, data)
    chart = graphic_frame.chart

    title = req.get("title") or ""
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    prs.save(path)
    return {
        "slide": slide_idx,
        "shape": len(slide.shapes),
        "type": chart_alias,
        "series": len(data.categories) if hasattr(data, "categories") else 0,
    }


def op_update(req: Dict[str, Any]) -> Dict[str, Any]:
    path = req["file"]
    prs = Presentation(path)
    slide_idx = int(req["slide"])
    shape_idx = int(req["shape"])
    if slide_idx < 1 or slide_idx > len(prs.slides):
        raise ValueError(f"slide {slide_idx} out of range")
    slide = prs.slides[slide_idx - 1]
    shapes = list(slide.shapes)
    if shape_idx < 1 or shape_idx > len(shapes):
        raise ValueError(f"shape {shape_idx} out of range")
    target = shapes[shape_idx - 1]
    if not target.has_chart:
        raise ValueError(f"shape {shape_idx} is not a chart")
    chart = target.chart

    # Replace data in place — keeps the chart's formatting / type / position.
    if chart.chart_type in (XL_CHART_TYPE.XY_SCATTER, XL_CHART_TYPE.XY_SCATTER_LINES,
                             XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, XL_CHART_TYPE.XY_SCATTER_SMOOTH,
                             XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS):
        data = _build_xy_data(req.get("csv", ""))
    else:
        data = _build_category_data(chart.chart_title.text_frame.text if chart.has_title else "", req.get("csv", ""))

    chart.replace_data(data)

    prs.save(path)
    return {"slide": slide_idx, "shape": shape_idx}


def op_delete(req: Dict[str, Any]) -> Dict[str, Any]:
    path = req["file"]
    prs = Presentation(path)
    slide_idx = int(req["slide"])
    shape_idx = int(req["shape"])
    if slide_idx < 1 or slide_idx > len(prs.slides):
        raise ValueError(f"slide {slide_idx} out of range")
    slide = prs.slides[slide_idx - 1]
    shapes = list(slide.shapes)
    if shape_idx < 1 or shape_idx > len(shapes):
        raise ValueError(f"shape {shape_idx} out of range")
    target = shapes[shape_idx - 1]
    if not target.has_chart:
        raise ValueError(f"shape {shape_idx} is not a chart")
    # Delete the underlying XML element (python-pptx exposes no remove API)
    sp = target._element  # noqa: SLF001 - intentional internal access
    sp.getparent().remove(sp)

    prs.save(path)
    return {"slide": slide_idx, "shape": shape_idx}


def op_charts(req: Dict[str, Any]) -> Dict[str, Any]:
    path = req["file"]
    prs = Presentation(path)
    return {"charts": _chart_list(prs, int(req["slide"]))}


_DISPATCH = {
    "charts": op_charts,
    "add":    op_add,
    "update": op_update,
    "delete": op_delete,
}


def main() -> None:
    if len(sys.argv) < 2:
        _fail("missing request JSON (argv[1])")
    try:
        req = json.loads(sys.argv[1])
    except Exception as e:
        _fail(f"invalid JSON: {e}")
        return
    op = req.get("op")
    handler = _DISPATCH.get(op or "")
    if not handler:
        _fail(f"unknown op: {op!r}")
        return
    try:
        result = handler(req)
    except Exception as e:
        tb = traceback.format_exc(limit=3)
        _fail(f"{e}\n{tb}")
        return
    sys.stdout.write(json.dumps({"ok": True, "result": result}))
    sys.stdout.flush()


if __name__ == "__main__":
    main()
